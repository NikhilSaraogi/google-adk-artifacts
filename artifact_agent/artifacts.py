"""
Artifact Management Tools — The engine behind the Artifact Agent.

This module provides the core logic for:
- Intercepting file uploads from the ADK native context.
- Converting Office documents (Excel, Word, PowerPoint) into text-based formats.
- Managing session-scoped artifact storage (save/load/list).
- Generating PDF reports from text content.
"""

import re
import mimetypes
import base64
import datetime
import csv
import os
from pathlib import Path
from io import BytesIO, StringIO
from typing import Dict, Any, Tuple, List, Optional

from google import genai
from google.genai import types
from google.adk.tools.tool_context import ToolContext

# ──────────────────────────────────────────────────────────────
# CONSTANTS & CATEGORIES
# ──────────────────────────────────────────────────────────────

OFFICE_EXTENSIONS = {".xlsx", ".xls", ".docx", ".doc", ".pptx", ".ppt"}
BINARY_EXTENSIONS = {".pdf", ".png", ".jpg", ".jpeg", ".gif", ".webp", ".zip", ".mp4", ".mp3"}
TEXT_EXTENSIONS = {".txt", ".csv", ".json", ".md", ".py", ".js", ".ts", ".sql", ".yaml", ".yml"}

# ──────────────────────────────────────────────────────────────
# INTERNAL HELPERS
# ──────────────────────────────────────────────────────────────

def _log(message: str):
    """Internal logging helper."""
    print(f"[Artifacts] {message}")


def _get_mime_type(filename: str) -> str:
    """Detect MIME type with custom overrides for Office/Code formats."""
    custom_mimes = {
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc":  "application/msword",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt":  "application/vnd.ms-powerpoint",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls":  "application/vnd.ms-excel",
        ".md":   "text/markdown",
        ".py": "text/plain", ".js": "text/plain", ".ts": "text/plain", ".sql": "text/plain",
    }
    ext = Path(filename).suffix.lower()
    if ext in custom_mimes: return custom_mimes[ext]
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or "application/octet-stream"


def _safe_base64_decode(data: Optional[str]) -> bytes:
    """Resilient base64 decoder with padding and cleanup."""
    if not isinstance(data, str) or not data: return b""
    # Handle data-URI prefix
    if "," in data: data = data.split(",", 1)[1]
    # Remove all whitespace and non-base64 chars
    data = "".join(data.split())
    data = re.sub(r'[^A-Za-z0-9+/=]', '', data)
    
    missing_padding = len(data) % 4
    if missing_padding: data += "=" * (4 - missing_padding)
    
    try:
        return base64.b64decode(data)
    except Exception as e:
        _log(f"Base64 decode failed: {e}")
        return b""


def _detect_file_format(raw_bytes: bytes) -> str:
    """Detect file format using magic bytes."""
    if not raw_bytes: return "empty"
    # ZIP magic (PK)
    if raw_bytes.startswith(b'PK\x03\x04'): return "zip"
    # OLE magic (legacy Office)
    if raw_bytes.startswith(b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1'): return "legacy"
    # PDF magic
    if raw_bytes.startswith(b'%PDF'): return "pdf"
    return "unknown"


def _convert_office_content(file_name: str, file_content: Optional[str], raw_bytes: Optional[bytes] = None) -> Tuple[bytes, str, str]:
    """
    Converts Office content to the requested format.
    Priority: Excel -> CSV, Word/PPT -> TXT.
    """
    ext = Path(file_name).suffix.lower()
    if raw_bytes is None:
        raw_bytes = _safe_base64_decode(file_content)
    
    if not raw_bytes:
        return b"Error: File content is empty or could not be decoded.", ".txt", "text/plain"

    fmt = _detect_file_format(raw_bytes)
    _log(f"Detected format for {file_name}: {fmt}")

    try:
        if ext in (".xlsx", ".xls"):
            import pandas as pd
            output = StringIO()
            try:
                # Explicitly specify engine
                engine = "openpyxl" if fmt == "zip" else "xlrd"
                dict_df = pd.read_excel(BytesIO(raw_bytes), sheet_name=None, engine=engine)
                for i, (sheet_name, df) in enumerate(dict_df.items()):
                    if i > 0: output.write("\n")
                    output.write(f"--- Sheet: {sheet_name} ---\n")
                    df.to_csv(output, index=False)
                return output.getvalue().encode("utf-8"), ".csv", "text/csv"
            except Exception as excel_err:
                _log(f"Excel conversion failed: {excel_err}")
                return f"Error: Failed to parse Excel file '{file_name}'. {excel_err}".encode("utf-8"), ".txt", "text/plain"
            
        elif ext in (".docx", ".doc"):
            if fmt == "legacy":
                return f"Error: '{file_name}' is in legacy binary format (.doc). Please save it as a modern .docx file and re-upload.".encode("utf-8"), ".txt", "text/plain"
            import docx
            try:
                doc = docx.Document(BytesIO(raw_bytes))
                text_content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                return text_content.encode("utf-8"), ".txt", "text/plain"
            except Exception as doc_err:
                return f"Error converting Word doc '{file_name}': {doc_err}".encode("utf-8"), ".txt", "text/plain"
            
        elif ext in (".pptx", ".ppt"):
            if fmt == "legacy":
                return f"Error: '{file_name}' is in legacy binary format (.ppt). Please save it as a modern .pptx file and re-upload.".encode("utf-8"), ".txt", "text/plain"
            from pptx import Presentation
            try:
                prs = Presentation(BytesIO(raw_bytes))
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slides_text.append(f"--- Slide {i} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slides_text.append(shape.text.strip())
                    slides_text.append("")
                return "\n".join(slides_text).encode("utf-8"), ".txt", "text/plain"
            except Exception as ppt_err:
                return f"Error converting PowerPoint '{file_name}': {ppt_err}".encode("utf-8"), ".txt", "text/plain"
            
        return raw_bytes, ext, "application/octet-stream"
    except Exception as e:
        _log(f"Unexpected error in _convert_office_content: {e}")
        return f"Unexpected error processing {file_name}: {e}".encode("utf-8"), ".txt", "text/plain"


async def _resolve_upload_content(tool_context: ToolContext, file_name: str, file_content: str, mime_type: str) -> Tuple[Optional[str], Optional[Any], Optional[bytes]]:
    """Scans ADK context and session state for upload content."""
    raw_data = None
    
    # 1. Check Native ADK Context (current turn)
    if hasattr(tool_context, "user_content") and tool_context.user_content and hasattr(tool_context.user_content, "parts"):
        for part in tool_context.user_content.parts:
            blob = part.inline_data or (part.file_data if hasattr(part, "file_data") else None)
            if not blob: continue
            
            p_mime = getattr(blob, "mime_type", None)
            p_name = getattr(blob, "display_name", None)
            
            # Match by name or MIME
            is_match = (p_name and file_name and p_name.lower() == file_name.lower()) or \
                       (p_mime and (p_mime == mime_type or p_mime.split('/')[0] == mime_type.split('/')[0] and p_mime.split('/')[0] not in ('application', 'octet-stream')))
            
            if is_match:
                _log(f"Matched '{file_name}' in native ADK context.")
                if hasattr(blob, "data") and blob.data:
                    raw_data = blob.data
                elif hasattr(blob, "file_uri") and blob.file_uri:
                    raw_data = await _read_from_uri(blob.file_uri)
                
                if raw_data and not file_content:
                    file_content = base64.b64encode(raw_data).decode('utf-8')
                return file_content, part, raw_data

    # 2. Check Session Cache (for stripped files)
    pending = tool_context.state.get("pending_office_uploads") if hasattr(tool_context, "state") else None
    if pending and isinstance(pending, list):
        for i, item in enumerate(pending):
            p_name = item.get("name")
            p_mime = item.get("mime_type")
            
            if (p_name and file_name and p_name.lower() == file_name.lower()) or (p_mime == mime_type):
                _log(f"Matched '{file_name}' in session cache.")
                if item.get("data"):
                    file_content = item["data"]
                    raw_data = _safe_base64_decode(file_content)
                elif item.get("file_uri"):
                    raw_data = await _read_from_uri(item["file_uri"])
                    if raw_data: file_content = base64.b64encode(raw_data).decode('utf-8')
                
                pending.pop(i)
                tool_context.state["pending_office_uploads"] = pending
                return file_content, None, raw_data
    
    if file_content: raw_data = _safe_base64_decode(file_content)
    return file_content, None, raw_data


async def _read_from_uri(uri: str) -> Optional[bytes]:
    """Helper to read bytes from a file:// URI or local path."""
    try:
        path = str(uri)
        if path.startswith("file://"): 
            path = path[7:]
        if os.path.exists(path) and os.path.isfile(path):
            _log(f"Reading local file: {path}")
            with open(path, "rb") as f:
                return f.read()
    except Exception as e:
        _log(f"Failed to read from URI {uri}: {e}")
    return None


async def _save_and_format_response(tool_context: ToolContext, file_name: str, data_bytes: bytes, mime_type: str, status_msg: str) -> dict:
    """Standardized helper to save artifact and return response."""
    try:
        artifact_part = types.Part(inline_data=types.Blob(mime_type=mime_type, data=data_bytes))
        await tool_context.save_artifact(file_name, artifact_part)
        return {
            "status": "success", "message": status_msg, "artifact_name": file_name,
            "mime_type": mime_type, "size_kb": round(float(len(data_bytes or b"")) / 1024.0, 2)
        }
    except Exception as e:
        return {"status": "error", "message": f"Failed to save '{file_name}': {e}"}


def _generate_pdf_blob(content: str) -> bytes:
    """Generates PDF blob from text content."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    pdf_buffer = BytesIO()
    doc = SimpleDocTemplate(pdf_buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    styles = getSampleStyleSheet()
    body_style = ParagraphStyle("Body", parent=styles["Normal"], fontSize=10, leading=14)
    elements = [Paragraph("Report", styles["Heading1"]), Spacer(1, 0.2*inch), Paragraph(f"<i>Generated: {datetime.datetime.now()}</i>", styles["Normal"]), Spacer(1, 0.3*inch)]
    for para in content.split("\n\n"):
        if para.strip():
            elements.append(Paragraph(para.replace("\n", " ").strip(), body_style))
            elements.append(Spacer(1, 0.1*inch))
    doc.build(elements)
    return pdf_buffer.getvalue()

# ──────────────────────────────────────────────────────────────
# PUBLIC TOOLS
# ──────────────────────────────────────────────────────────────

async def handle_file_upload(tool_context: ToolContext, file_name: str, file_content: str = "") -> dict:
    """
    Primary entry point for saving uploads. 
    
    This function automatically:
    1. Resolves the raw file content from either the direct 'file_content' argument, 
       the current ADK turn's user content, or the session-scoped cache.
    2. Detects the file type and applies necessary conversions (e.g., Office to CSV).
    3. Persists the final data as a session artifact.
    
    Args:
        tool_context: The ADK ToolContext providing access to state and artifact services.
        file_name: The original name of the file (e.g., 'data.xlsx').
        file_content: Optional base64-encoded string (fallback if not in context).
        
    Returns:
        dict: { status: "success"|"error", message: str, artifact_name: str, ... }
    """
    try:
        if not file_name: return {"status": "error", "message": "file_name is required."}
        
        ext = Path(file_name).suffix.lower()
        mime_type = _get_mime_type(file_name)
        _log(f"🚀 Processing upload: '{file_name}' (MIME: {mime_type})")
        
        file_content, found_raw_part, raw_bytes = await _resolve_upload_content(tool_context, file_name, file_content, mime_type)

        if not raw_bytes and not file_content and not found_raw_part:
            _log(f"ERROR: No content found for '{file_name}'.")
            return {"status": "error", "message": f"No content for '{file_name}'. Ensure the file was uploaded correctly."}

        # ✨ ALWAYS handle Office files with specific format logic
        if ext in OFFICE_EXTENSIONS:
            _log(f"Handling Office file: {file_name}")
            data_bytes, new_ext, new_mime = _convert_office_content(file_name, file_content, raw_bytes)
            new_file_name = file_name.rsplit(".", 1)[0] + new_ext
            status_msg = f"✅ '{file_name}' processed and saved as '{new_file_name}'."
            return await _save_and_format_response(tool_context, new_file_name, data_bytes, new_mime, status_msg)

        if found_raw_part:
            _log(f"Saving '{file_name}' natively from artifact part.")
            await tool_context.save_artifact(file_name, found_raw_part)
            return {"status": "success", "message": f"Saved '{file_name}' natively.", "artifact_name": file_name}

        # ✨ Standard Binary / Multimedia handling
        if ext in BINARY_EXTENSIONS or mime_type.startswith(("image/", "video/", "audio/", "application/pdf")):
            _log(f"Saving '{file_name}' as binary artifact.")
            if not raw_bytes and file_content:
                raw_bytes = _safe_base64_decode(file_content)
            
            if not raw_bytes:
                return {"status": "error", "message": f"Failed to resolve or decode binary data for '{file_name}'."}
                
            return await _save_and_format_response(tool_context, file_name, raw_bytes, mime_type, f"✅ Binary '{file_name}' saved.")

        return await _save_and_format_response(tool_context, file_name, (file_content or "").encode("utf-8"), mime_type if ext in TEXT_EXTENSIONS else "text/plain", f"✅ File '{file_name}' saved.")
    except Exception as e:
        return {"status": "error", "message": f"Unexpected error: {e}"}


async def load_artifacts(tool_context: ToolContext, artifact_name: str) -> dict:
    """
    Loads the content of a previously stored artifact.
    
    Args:
        tool_context: ADK ToolContext.
        artifact_name: The name of the artifact to retrieve.
        
    Returns:
        dict: { status: "success", artifact_name: str, artifact: types.Part }
    """
    try:
        artifact = await tool_context.load_artifact(artifact_name)
        if artifact is None:
            keys = await tool_context.list_artifact_keys()
            _log(f"❌ Artifact '{artifact_name}' not found.")
            return {"status": "error", "message": f"Artifact '{artifact_name}' not found. Available: {keys}"}
        
        _log(f"✅ Loaded artifact: '{artifact_name}'")
        
        # Extract text content for easy model reading
        content_preview = None
        mime_type = "unknown"
        
        blob = getattr(artifact, "inline_data", None) or getattr(artifact, "file_data", None)
        if blob:
            mime_type = getattr(blob, "mime_type", "unknown")
            if mime_type.startswith("text/") or mime_type in ("application/json", "text/csv"):
                if hasattr(blob, "data") and blob.data:
                    try:
                        content_preview = blob.data.decode("utf-8")
                    except:
                        content_preview = "[Binary or unencodable text]"
        
        response = {
            "status": "success", 
            "artifact_name": artifact_name, 
            "mime_type": mime_type,
            "artifact": artifact
        }
        if content_preview:
            response["content_preview"] = content_preview
            
        return response
    except Exception as e:
        _log(f"❌ Error loading '{artifact_name}': {e}")
        return {"status": "error", "message": str(e)}


async def list_artifact_keys(tool_context: ToolContext) -> dict:
    """List all artifact names saved in the current session."""
    try:
        keys = await tool_context.list_artifact_keys()
        return {"status": "success", "artifact_keys": keys or [], "count": len(keys) if keys else 0}
    except Exception as e:
        return {"status": "error", "message": str(e)}




async def extract_data_from_artifact(
    tool_context: ToolContext,
    artifact_name: str,
    specific_request: str
) -> str:
    """
    Directly extracts specific information or data from an artifact using native Gemini vision/processing.
    Use this for technical details, formulas, or descriptions that require high precision.
    """
    try:
        artifact_part = await tool_context.load_artifact(artifact_name)
        if not artifact_part:
            return f"Error: Artifact '{artifact_name}' not found."
            
        _log(f"Deep extraction from '{artifact_name}' for: {specific_request}")
        
        from google import genai
        # Rebuild standard GenAI Part to ensure compatibility after ADK deserialization
        genai_part = artifact_part
        if isinstance(artifact_part, dict):
            # Parse from dictionary
            inline_data = artifact_part.get("inline_data") or artifact_part.get("inlineData")
            if inline_data:
                mime_type = inline_data.get("mime_type") or inline_data.get("mimeType")
                data = inline_data.get("data")
                if isinstance(data, str):
                    data = base64.b64decode(data)
                genai_part = types.Part.from_bytes(data=data, mime_type=mime_type)
        else:
            # Maybe it's an object with inline_data
            inline_data = getattr(artifact_part, "inline_data", None)
            if inline_data:
                data = getattr(inline_data, "data", None)
                mime = getattr(inline_data, "mime_type", None)
                if isinstance(data, str):
                     data = base64.b64decode(data)
                if data and mime:
                     genai_part = types.Part.from_bytes(data=data, mime_type=mime)
        
        prompt = f'''
Please review this document and extract the following specific information:
"{specific_request}"

CRITICAL RULES FOR EXTRACTION:
1. ONLY extract exactly what is explicitly written in the document.
2. DO NOT guess, infer, or invent ANY part of the formula, variables, units, or descriptions.
3. If a value or parameter is missing from the document, LEAVE IT BLANK.
4. If the information is NOT found, reply with: "Error: The requested information was not found in the uploaded document."

Return the results in a structured format (Tables/Markdown).
'''
        from .utils import model
        model_name = getattr(model, "model", "gemini-2.0-flash")
        
        client = genai.Client()
        response = client.models.generate_content(
            model=model_name,
            contents=[genai_part, prompt]
        )
        return response.text
    except Exception as e:
        return f"Failed to extract from artifact: {str(e)}"


async def save_artifact_content(tool_context: ToolContext, content: str, filename: str, output_format: str = "text") -> dict:
    """Save text content (reports, analysis) as a session artifact (.txt, .md, .pdf)."""
    if not content: return {"status": "error", "message": "Content empty."}
    output_format = output_format.lower().strip()
    mime_type = {"markdown": "text/markdown", "pdf": "application/pdf"}.get(output_format, "text/plain")

    try:
        if output_format == "pdf":
            try:
                data_bytes = _generate_pdf_blob(content)
                return await _save_and_format_response(tool_context, filename, data_bytes, mime_type, f"✅ PDF '{filename}' saved.")
            except Exception:
                filename = filename.rsplit(".", 1)[0] + ".txt"
                mime_type = "text/plain"
        return await _save_and_format_response(tool_context, filename, content.encode("utf-8"), mime_type, f"✅ {output_format.upper()} '{filename}' saved.")
    except Exception as e:
        return {"status": "error", "message": str(e)}